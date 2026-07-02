"""PPTX converter — delegates to MarkItDown with python-pptx fallback."""
import os
import re
from typing import Any, BinaryIO, List

from core.harness.document.protocol import (
    DocumentConverter, DocumentElement, StreamInfo, detect_structure_role,
)


class PptxConverter(DocumentConverter):
    """PPTX → Markdown via MarkItDown (or python-pptx fallback)."""

    SOURCE_FORMAT = "pptx"
    REQUIRED_PACKAGES = {}  # markitdown + python-pptx are soft deps
    ACCEPTED_EXTENSIONS = (".pptx", ".ppt")
    ACCEPTED_MIME_PREFIXES = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    )

    def accepts(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> bool:
        return self._accepts_by_format(stream_info)

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,
    ) -> List[DocumentElement]:
        file_path = stream_info.local_path
        if not file_path:
            return self._fallback_text("(stream)", stream_info)

        try:
            from markitdown import MarkItDown
            return self._convert_via_markitdown(file_path)
        except ImportError:
            return self._fallback_pptx(file_path)

    def _fallback_pptx(self, file_path: str) -> List[DocumentElement]:
        try:
            from pptx import Presentation
        except ImportError:
            return self._fallback_text(file_path, StreamInfo(extension=".pptx"))

        prs = Presentation(file_path)
        elements: List[DocumentElement] = []
        for si, slide in enumerate(prs.slides):
            texts = [
                para.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame
                for para in shape.text_frame.paragraphs
                if (para.text or "").strip()
            ]
            if texts:
                elements.append(DocumentElement(
                    type="text", text="\n".join(texts),
                    page_idx=si,
                    meta={"source": "pptx", "slide_index": si},
                    source_format="pptx",
                ))
        return elements or self._fallback_text(file_path, StreamInfo(extension=".pptx"))

    def _fallback_text(self, file_path: str, stream_info: StreamInfo) -> List[DocumentElement]:
        if os.path.isfile(file_path):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            return []
        if not text.strip():
            return []
        return [DocumentElement(
            type="text", text=text.strip(), page_idx=0,
            meta={"source": "pptx", "fallback": True},
            source_format="pptx",
        )]
