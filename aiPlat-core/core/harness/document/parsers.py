"""
Core document parsers — unified document format → text elements.

Each parser takes a file path and returns a list of element dicts:
  {type: "text"|"table", text: str, page_idx: int, cells: Optional[List], meta: dict}

All parser functions are thin wrappers that delegate to the ConverterRegistry.
The registry is the SINGLE SOURCE OF TRUTH for format→converter dispatch.

Callers:
  - platform/kb/service.py (ingest pipeline)
  - Any AI agent that needs to read documents before reasoning
"""
from __future__ import annotations
import logging

import os
import re
from typing import Any, Dict, List

from core.harness.document.protocol import StreamInfo


def _elements_to_dicts(elements: list) -> List[Dict[str, Any]]:
    """Convert DocumentElement list to legacy dict format for backward compatibility."""
    result = []
    for el in elements:
        result.append({
            "type": el.type,
            "text": el.text,
            "page_idx": el.page_idx,
            "cells": el.cells,
            "meta": el.meta,
        })
    return result


def _resolve_parser(file_path: str, formats: List[str]):
    """Find the right converter for a given file path and format list."""
    from core.harness.document.protocol import get_document_registry, StreamInfo
    ext = os.path.splitext(file_path)[1].lower()
    registry = get_document_registry()
    info = StreamInfo(local_path=file_path, extension=ext)
    for fmt in formats:
        converter = registry.find_converter(StreamInfo(extension=f".{fmt}"))
        if converter:
            return converter, info
        converter = registry.find_converter(StreamInfo(extension=ext))
        if converter:
            return converter, info
    return None, info


# ── DOCX ──

def parse_docx(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["docx", "doc"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── PPTX ──

def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["pptx", "ppt"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── MARKDOWN ──

def parse_markdown(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["md", "markdown"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── CSV ──

def parse_csv(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["csv"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── MarkItDown (unified DOCX/PPTX/XLSX/PDF/HTML → structured Markdown) ──

def parse_markitdown(file_path: str) -> List[Dict[str, Any]]:
    ext = os.path.splitext(file_path)[1].lower()
    registry = get_document_registry()
    info = StreamInfo(local_path=file_path, extension=ext)
    converter = registry.find_converter(info)
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))

    ext_no_dot = ext.lstrip(".")
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if not text.strip():
        return []
    return [{"type": "text", "text": text.strip(), "page_idx": 0, "cells": None,
             "meta": {"source": f"markitdown+{ext_no_dot}", "fallback": True}}]


# ── HTML ──

def parse_html(file_path: str) -> List[Dict[str, Any]]:
    return parse_markitdown(file_path)


# ── Audio ──

def parse_audio(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["audio", "mp3", "wav", "m4a"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── Image ──

def parse_image(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["image", "png", "jpg", "jpeg"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── JSON ──

def parse_json_document(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["json"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


# ── Email .eml ──

def parse_eml(file_path: str) -> List[Dict[str, Any]]:
    converter, info = _resolve_parser(file_path, ["eml"])
    if converter:
        with open(file_path, "rb") as f:
            return _elements_to_dicts(converter.convert(f, info))
    return []


from core.harness.document.protocol import get_document_registry

__all__ = ["parse_docx", "parse_pptx", "parse_markdown", "parse_csv", "parse_audio", "parse_image", "parse_json_document", "parse_eml", "parse_markitdown", "parse_html", "extract_images_from_document", "describe_images"]


# ── Image extraction from documents ──

def extract_images_from_document(file_path: str, output_dir: str = None) -> List[str]:
    """Extract embedded images from DOCX/PDF/PPTX to a temp directory."""
    import tempfile
    ext = os.path.splitext(file_path)[1].lower()
    out = output_dir or tempfile.mkdtemp(prefix="aiplat_img_")
    os.makedirs(out, exist_ok=True)
    extracted: List[str] = []

    if ext in (".docx", ".doc"):
        extracted = _extract_docx_images(file_path, out)
    elif ext == ".pdf":
        extracted = _extract_pdf_images(file_path, out)
    elif ext in (".pptx", ".ppt"):
        extracted = _extract_pptx_images(file_path, out)
    return extracted


def _extract_docx_images(file_path: str, out_dir: str) -> List[str]:
    extracted = []
    try:
        from docx import Document
        doc = Document(file_path)
        for i, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                img_data = rel.target_part.blob
                ext2 = os.path.splitext(rel.target_part.partname)[1] or ".png"
                img_path = os.path.join(out_dir, f"docx_img_{i}{ext2}")
                with open(img_path, "wb") as f:
                    f.write(img_data)
                extracted.append(img_path)
    except ImportError:
        pass  # noqa: optional-dependency
    return extracted


def _extract_pdf_images(file_path: str, out_dir: str) -> List[str]:
    extracted = []
    try:
        import fitz
        doc = fitz.open(file_path)
        for pi in range(len(doc)):
            for img_info in doc[pi].get_images(full=True):
                xref = img_info[0]
                base_img = doc.extract_image(xref)
                img_bytes = base_img["image"]
                ext2 = base_img.get("ext", "png")
                img_path = os.path.join(out_dir, f"pdf_p{pi}_img_{xref}.{ext2}")
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                extracted.append(img_path)
        doc.close()
    except ImportError:
        pass  # noqa: optional-dependency
    return extracted


def _extract_pptx_images(file_path: str, out_dir: str) -> List[str]:
    extracted = []
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
        for si, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shape.image
                    ct = img.content_type or "image/png"
                    ext2 = ct.split("/")[-1] if "/" in ct else "png"
                    img_path = os.path.join(out_dir, f"pptx_s{si}_{shape.shape_id}.{ext2}")
                    with open(img_path, "wb") as f:
                        f.write(img.blob)
                    extracted.append(img_path)
    except ImportError:
        pass  # noqa: optional-dependency
    return extracted


# ── Image description via infra LLM vision ──

async def describe_images(image_paths: List[str],
                           model_name: str = None,
                           prompt: str = "请用2-3句中文描述这张图片的主要内容，包括图中展示的结构、流程或关键信息。") -> List[Dict[str, str]]:
    """Generate semantic descriptions for images using a vision-capable LLM."""
    results = []
    if not image_paths:
        return results

    import base64

    try:
        from core.harness.utils.model_injection import create_selected_adapter, best_model_for_purpose, get_default_model
        if not model_name:
            model_name = best_model_for_purpose("vision") or best_model_for_purpose("multimodal") or get_default_model(purpose="vision")
        if not model_name:
            return [{"path": p, "description": "", "error": "no vision model configured"}
                    for p in image_paths]
        adapter = create_selected_adapter(model_name=model_name)
    except Exception:
        return [{"path": p, "description": "", "error": "model init failed"}
                for p in image_paths]

    for img_path in image_paths[:8]:
        try:
            with open(img_path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode("utf-8")
            ext3 = os.path.splitext(img_path)[1].lower().strip(".")
            mime = f"image/{ext3}" if ext3 in ("png", "jpeg", "jpg", "gif", "webp") else "image/png"
            messages = [{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{img_b64}"}},
                ]
            }]
            resp = await adapter.generate(messages, config=None)
            desc = resp.content if hasattr(resp, 'content') else str(resp)
            results.append({"path": img_path, "description": desc.strip(), "error": None})
        except Exception as e:
            results.append({"path": img_path, "description": "", "error": str(e)[:200]})

    return results

